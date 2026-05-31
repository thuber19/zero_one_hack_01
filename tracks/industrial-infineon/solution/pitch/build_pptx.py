#!/usr/bin/env python3
"""
Build an EDITABLE 16:9 PowerPoint pitch deck for the ProcSeq hackathon submission.

Mirrors pitch/index.html (Infineon dark-tech theme, 7 slides) using native
PowerPoint primitives — text boxes, rounded rectangles, an arrow flow, and a
real table — so every element is editable in PowerPoint. Speaker notes included.

Run:  .pptxvenv/bin/python build_pptx.py
Out:  ProcSeq_Pitch.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- theme ---------------------------------------------------------------
BG      = RGBColor.from_string("0B1220")
BG2     = RGBColor.from_string("0E1626")
SURFACE = RGBColor.from_string("141E32")
TEXT    = RGBColor.from_string("E8EDF5")
MUTED   = RGBColor.from_string("9FB0C8")
FAINT   = RGBColor.from_string("6B7C99")
ACCENT  = RGBColor.from_string("EC0016")   # Infineon red
ACCENT2 = RGBColor.from_string("FF5A45")
BLUE    = RGBColor.from_string("3678FF")
LINE    = RGBColor.from_string("243450")
REDPILL = RGBColor.from_string("FFB3AD")

SANS = "Arial"
MONO = "Consolas"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ---- helpers -------------------------------------------------------------
def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def _noline(shape):
    shape.line.fill.background()


def rect(s, x, y, w, h, fill=SURFACE, line=LINE, lw=1.0, rounded=True, radius=0.10):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        _noline(shp)
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    shp.shadow.inherit = False
    return shp


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    return tb, tf


def para(tf, segments, size=18, color=MUTED, bold=False, font=SANS,
         align=PP_ALIGN.LEFT, space_after=6, leading=1.12, first=False):
    """segments: str or list of (text, {overrides})."""
    if first and not tf.paragraphs[0].runs:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    try:
        p.line_spacing = leading
    except Exception:
        pass
    if isinstance(segments, str):
        segments = [(segments, {})]
    for txt, ov in segments:
        r = p.add_run(); r.text = txt
        f = r.font
        f.size = Pt(ov.get("size", size))
        f.bold = ov.get("bold", bold)
        f.name = ov.get("font", font)
        f.color.rgb = ov.get("color", color)
    return p


def kicker(s, text):
    rect(s, 0.9, 0.64, 0.34, 0.05, fill=ACCENT, line=None, rounded=False)
    _, tf = textbox(s, 1.34, 0.5, 10.5, 0.4)
    para(tf, [(text.upper(), {})], size=12.5, color=ACCENT2, bold=True, font=MONO, first=True)


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def card(s, x, y, w, h, title_segs, body_segs, accent=LINE):
    rect(s, x, y, w, h, fill=SURFACE, line=accent, lw=1.25, radius=0.08)
    _, tf = textbox(s, x + 0.28, y + 0.24, w - 0.56, h - 0.48)
    para(tf, title_segs, size=17, color=TEXT, bold=True, first=True, space_after=8)
    para(tf, body_segs, size=12.5, color=MUTED, leading=1.22, space_after=0)


# =========================================================================
# 1 · TITLE
# =========================================================================
s = slide()
kicker(s, "Zero-One Hack 2026 · Infineon · Industrial AI")
_, tf = textbox(s, 0.9, 1.7, 11.5, 2.0)
para(tf, [("Process Logic, ", {"color": TEXT}), ("Learned.", {"color": ACCENT2})],
     size=58, bold=True, first=True, leading=1.02, space_after=14)
_, tf = textbox(s, 0.9, 3.6, 8.6, 1.2)
para(tf, "A hybrid neural–symbolic pipeline that learns the hidden grammar of "
         "semiconductor manufacturing — and proves it.",
     size=20, color=MUTED, leading=1.3, first=True)
pills = [("Next-step", True), ("Completion", True), ("Anomaly detection", True),
         ("+ hidden 4th-family OOD", False)]
px = 0.9
for label, red in pills:
    w = 0.5 + 0.108 * len(label)
    rect(s, px, 4.95, w, 0.46, fill=BG2, line=(ACCENT if red else LINE), lw=1.1, radius=0.5)
    _, f = textbox(s, px, 4.95, w, 0.46, anchor=MSO_ANCHOR.MIDDLE)
    para(f, [(label, {})], size=12, color=(REDPILL if red else MUTED),
         font=MONO, align=PP_ALIGN.CENTER, first=True)
    px += w + 0.2
_, tf = textbox(s, 0.9, 5.95, 11, 1.0)
para(tf, [("Team TBD", {})], size=16, color=TEXT, bold=True, first=True, space_after=4)
para(tf, [("Tobias Huber · Mina Mikail · Khaled El Yamany · Fathy Shalaby", {})],
     size=12.5, color=FAINT, font=MONO)
notes(s, "30s. Hook: a chip recipe is a long ordered sequence; we taught a model its "
         "grammar and built a hybrid that guarantees valid output. Three tasks plus a "
         "hidden out-of-distribution family. Team ProcSeq.")

# =========================================================================
# 2 · PROBLEM
# =========================================================================
s = slide()
kicker(s, "The problem")
_, tf = textbox(s, 0.9, 1.25, 11.6, 1.3)
para(tf, [("A chip is a ", {"color": TEXT}), ("150-step sentence", {"color": ACCENT2}),
          (" where order is physics.", {"color": TEXT})],
     size=33, bold=True, first=True, leading=1.06)
bullets = [
    [("Every wafer runs an ordered route: ", {}),
     ("clean → oxidize → pattern → etch → implant → … → test → ship", {"color": TEXT, "bold": True}),
     ("  (~115–150 steps, 3 families).", {})],
    [("Order encodes real constraints: ", {}),
     ("no etch before a mask, no deposit on a dirty surface, no ship before sort-test", {"color": TEXT, "bold": True}),
     (" — 10 long-range rules.", {})],
    [("The judges' real question: does the model learn the process ", {}),
     ("logic", {"color": ACCENT2, "bold": True}),
     (", or just memorize patterns?", {})],
    [("And it's graded on a ", {}),
     ("hidden 4th family", {"color": TEXT, "bold": True}),
     (" never seen in training.", {})],
]
_, tf = textbox(s, 0.9, 3.0, 11.5, 3.6)
for i, b in enumerate(bullets):
    para(tf, [("▸  ", {"color": ACCENT, "font": MONO})] + b,
         size=18, color=MUTED, leading=1.3, space_after=16, first=(i == 0))
notes(s, "40s. Process routes are long and order-sensitive — 10 hard rules, some spanning "
         "15 steps. The track explicitly asks whether a model LEARNS the logic or memorizes, "
         "and grades generalization on an unseen 4th product family.")

# =========================================================================
# 3 · INSIGHT
# =========================================================================
s = slide()
kicker(s, "The insight")
_, tf = textbox(s, 0.9, 1.25, 11.6, 1.3)
para(tf, [("Treat a recipe like a ", {"color": TEXT}), ("sentence", {"color": ACCENT2}),
          (". One step = one word.", {"color": TEXT})],
     size=33, bold=True, first=True, leading=1.06)
cw, cg, cy, ch = 3.71, 0.30, 2.85, 2.55
card(s, 0.9, cy, cw, ch,
     [("TOKENIZER", {"color": ACCENT2, "font": MONO, "size": 12})],
     [("A custom atomic-step tokenizer — each whole step is one token (vocab 210). "
       "No sub-word fragmentation; recipes read and write like language.", {})], accent=LINE)
card(s, 0.9 + (cw + cg), cy, cw, ch,
     [("The Writer · ", {"color": TEXT}), ("DECODER", {"color": ACCENT2, "font": MONO, "size": 12})],
     [("A from-scratch Llama-style causal Transformer that generates — next step & "
       "completion (Tasks 1–2). Attention does long-range lookups an LSTM can't.", {})], accent=BLUE)
card(s, 0.9 + 2 * (cw + cg), cy, cw, ch,
     [("The Judge · ", {"color": TEXT}), ("ENCODER", {"color": ACCENT2, "font": MONO, "size": 12})],
     [("A from-scratch DeBERTa-v2 encoder that classifies validity + which rule broke "
       "(Task 3), trained with contrastive hard-negative twins.", {})], accent=ACCENT)
_, tf = textbox(s, 0.9, 5.7, 11.5, 1.2)
para(tf, "Two specialized Transformers, trained from scratch — no pretrained model, full "
         "control of the vocabulary, interpretable representations.",
     size=16, color=MUTED, leading=1.3, first=True)
notes(s, "30s. The trick: one process step = one token (vocab 210). Two from-scratch "
         "Transformers — a decoder that writes routes and a DeBERTa encoder that judges them, "
         "trained on valid-vs-broken twins so it learns WHY a route is invalid.")

# =========================================================================
# 4 · HYBRID
# =========================================================================
s = slide()
kicker(s, "Our approach · the hybrid")
_, tf = textbox(s, 0.9, 1.25, 11.6, 1.3)
para(tf, [("The model proposes. ", {"color": TEXT}), ("Physics disposes.", {"color": ACCENT2})],
     size=33, bold=True, first=True, leading=1.06)
bw, bh, by = 3.45, 2.5, 3.0
bx = 0.9
boxdata = [
    ("NEURAL — PLAUSIBILITY", "Decoder + Encoder",
     "Learn the grammar from data. Rank likely next steps, generate fluent completions, "
     "and output a real probability of validity.", BLUE),
    ("SYMBOLIC — CORRECTNESS", "Physics Refinery",
     "A deterministic rule engine (≡ the official validator). Vetoes illegal steps, "
     "repairs dead-ends, certifies the verdict.", ACCENT),
    ("OUTPUT", "Best of both",
     "Guaranteed-valid routes that are also likely — “the model decided, the rules "
     "certify it.”", LINE),
]
positions = [bx, bx + bw + 0.74, bx + 2 * (bw + 0.74)]
for (lbl, title, body, acc), xpos in zip(boxdata, positions):
    rect(s, xpos, by, bw, bh, fill=SURFACE, line=acc, lw=1.4, radius=0.07)
    _, f = textbox(s, xpos + 0.26, by + 0.22, bw - 0.52, bh - 0.44)
    para(f, [(lbl, {})], size=11, color=FAINT, font=MONO, bold=True, first=True, space_after=8)
    para(f, [(title, {})], size=17, color=TEXT, bold=True, space_after=8)
    para(f, [(body, {})], size=12.5, color=MUTED, leading=1.22)
for ax in (bx + bw + 0.1, bx + 2 * bw + 0.84):
    ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(ax), Inches(by + bh / 2 - 0.22),
                            Inches(0.54), Inches(0.44))
    ar.fill.solid(); ar.fill.fore_color.rgb = ACCENT
    _noline(ar); ar.shadow.inherit = False
_, tf = textbox(s, 0.9, 5.95, 11.6, 1.1)
para(tf, [("Neural alone can be ", {}), ("confidently wrong", {"color": ACCENT2}),
          ("; rules alone can't say what's ", {}), ("likely", {"color": ACCENT2}),
          (". The hybrid covers both blind spots.", {})],
     size=16.5, color=TEXT, font=MONO, leading=1.3, first=True)
notes(s, "40s — KEY SLIDE. 'Model proposes, physics disposes.' Learned models supply "
         "plausibility; a deterministic rule engine (identical to the official validator) "
         "supplies correctness. Neural can be confidently wrong; rules can't rank. Together "
         "you get valid AND likely outputs.")

# =========================================================================
# 5 · THE 3 TASKS
# =========================================================================
s = slide()
kicker(s, "How we tackled the 3 tasks")
_, tf = textbox(s, 0.9, 1.25, 11.6, 1.3)
para(tf, [("One hybrid, wired ", {"color": TEXT}), ("three ways.", {"color": ACCENT2})],
     size=33, bold=True, first=True, leading=1.06)
tasks = [
    ([("1", {"color": ACCENT2, "font": MONO}), ("   Next-step", {"color": TEXT})],
     "Decoder's top-5 candidates, then the Refinery floats the legal ones to the top. "
     "Cannot hurt a good model; rescues a bad rank."),
    ([("2", {"color": ACCENT2, "font": MONO}), ("   Completion", {"color": TEXT})],
     "Decoder proposes; rule-vetoed beam search + repair (beam 5, branch 8). Every "
     "emitted suffix is provably valid."),
    ([("3", {"color": ACCENT2, "font": MONO}), ("   Anomaly", {"color": TEXT})],
     "Contrastive encoder learns WHY a route is invalid; rule engine gives the verdict, "
     "the model gives the graded score (real ROC-AUC)."),
]
cw, cy, ch = 3.71, 2.85, 2.55
for i, (title, body) in enumerate(tasks):
    card(s, 0.9 + i * (cw + 0.30), cy, cw, ch, title, [(body, {})], accent=LINE)
_, tf = textbox(s, 0.9, 5.7, 11.5, 1.2)
para(tf, "Generative tasks use the writer; the discriminative task uses the judge — each "
         "with a guarantee the learned model can't make alone.",
     size=16, color=MUTED, leading=1.3, first=True)
notes(s, "35s. Same two components, three wirings. Task 1: top-5 + legal-first re-sort. "
         "Task 2: rule-vetoed beam search with repair → guaranteed-valid completion. Task 3: "
         "contrastive encoder for a graded score, rule engine for the hard verdict.")

# =========================================================================
# 6 · RESULTS
# =========================================================================
s = slide()
kicker(s, "Results · official scorer, held-out split")
_, tf = textbox(s, 0.9, 1.2, 11.6, 1.2)
para(tf, [("It works — and we can ", {"color": TEXT}), ("prove it learned", {"color": ACCENT2}),
          (".", {"color": TEXT})], size=31, bold=True, first=True, leading=1.06)

rows = [
    ("Task", "Headline", "Also"),
    ("1 · Next-step", "Top-5 1.000  ·  Top-1 0.772", "MRR 0.88"),
    ("2 · Completion", "Block-level 0.92  ·  100% rule-valid", "token 0.57 · exact 0.14"),
    ("3 · Anomaly", "Hybrid verdict — exact in-dist", "encoder alone AUC 0.61 (honest)"),
    ("Understanding", "0.963 next-operation acc", "learns function, not just names"),
]
gt = s.shapes.add_table(len(rows), 3, Inches(0.9), Inches(2.7),
                        Inches(11.5), Inches(2.55)).table
gt.columns[0].width = Inches(3.0)
gt.columns[1].width = Inches(5.0)
gt.columns[2].width = Inches(3.5)
gt._tbl.tblPr.set('firstRow', '0'); gt._tbl.tblPr.set('bandRow', '0')
for r in range(len(rows)):
    for c in range(3):
        cell = gt.cell(r, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BG2 if r == 0 else (SURFACE if r % 2 else BG)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.16); cell.margin_right = Inches(0.1)
        cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
        tfc = cell.text_frame; tfc.word_wrap = True
        p = tfc.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        run = p.add_run(); run.text = rows[r][c]; f = run.font
        if r == 0:
            f.size = Pt(11); f.bold = True; f.name = MONO; f.color.rgb = FAINT
        elif c == 0:
            f.size = Pt(14); f.bold = True; f.name = SANS; f.color.rgb = TEXT
        elif c == 1:
            f.size = Pt(15); f.bold = True; f.name = MONO; f.color.rgb = ACCENT2
        else:
            f.size = Pt(12.5); f.name = MONO; f.color.rgb = MUTED

_, tf = textbox(s, 0.9, 5.45, 11.6, 0.55)
para(tf, [("learns ", {}), ("function", {"color": ACCENT2, "bold": True}),
          (", not just names: predicts the right ", {}),
          ("operation 96% of the time", {"color": ACCENT2, "bold": True}),
          ("   ·   completions are 100% rule-valid", {})],
     size=13, color=MUTED, font=MONO, first=True)
_, tf = textbox(s, 0.9, 6.15, 11.6, 1.0)
para(tf, [("Honest provenance: ", {"color": MUTED, "bold": True}),
          ("real numbers from the organizer's official scorer on a held-out self-eval split, "
           "from the completed ProcSeq Leonardo run (base, 16k steps). The learned anomaly "
           "encoder alone is weak (AUC 0.61) — the physics hybrid carries Task 3. Organizer "
           "hidden-test and 4th-family OOD are pending.", {"color": FAINT})],
     size=10.5, leading=1.3, first=True)
notes(s, "35s. Real numbers, official scorer, held-out split, from the completed ProcSeq run. "
         "Task1: Top-5 perfect, Top-1 0.77, and 0.96 next-OPERATION accuracy → it learned function, "
         "not just names. Task2: block-level 0.92 and 100% rule-valid completions. Task3, be honest: "
         "the learned encoder alone is ~chance (AUC 0.61), so the physics hybrid gives the exact "
         "in-distribution verdict. Organizer hidden test + 4th-family OOD still pending.")

# =========================================================================
# 7 · CLOSE
# =========================================================================
s = slide()
kicker(s, "Why this submission")
_, tf = textbox(s, 0.9, 1.35, 11.0, 1.7)
para(tf, [("We measured ", {"color": TEXT}), ("whether the model learned", {"color": ACCENT2}),
          (" — not just that it scored.", {"color": TEXT})],
     size=30, bold=True, first=True, leading=1.12)
closecards = [
    ("Honest evaluation",
     "Every number sits between an n-gram floor and a rule-oracle ceiling, plus a logic "
     "probe & a novel-vocabulary OOD stress test we built to make it fail."),
    ("Hybrid & robust",
     "Learned plausibility + symbolic guarantees → 0 illegal outputs, a calibratable "
     "anomaly score, and a graceful OOD story."),
    ("Reproducible",
     "make smoke in <1 min; one command trains the full pipeline on Leonardo. Unit tests, "
     "fixed seeds, fully open stack."),
]
cw, cy, ch = 3.71, 3.2, 2.4
for i, (title, body) in enumerate(closecards):
    card(s, 0.9 + i * (cw + 0.30), cy, cw, ch, [(title, {})], [(body, {})], accent=LINE)
_, tf = textbox(s, 0.9, 5.95, 11.6, 1.0)
para(tf, [("ProcSeq", {"color": TEXT, "bold": True}),
          (" — process logic, learned, guaranteed, and honestly benchmarked.   ", {"color": MUTED}),
          ("Thank you.", {"color": ACCENT2, "bold": True})],
     size=17, leading=1.3, first=True)
notes(s, "20s close. Our differentiator is honesty: floor-to-ceiling ladder, logic probe, "
         "and a novel-vocab OOD test we designed to break our own model. Hybrid gives 0 "
         "illegal outputs. Fully reproducible on Leonardo. Thank you — questions?")

# ---- save ----------------------------------------------------------------
OUT = "ProcSeq_Pitch.pptx"
prs.save(OUT)
print(f"Wrote {OUT} — {len(prs.slides._sldIdLst)} slides")
